"""Local document library and per-document PixelRAG lifecycle for the POC.

The production StratAlign equivalent would persist metadata in PostgreSQL, source
files in object storage, and indexing work in a job queue.  This POC keeps the
same boundaries while using the local filesystem so the feature stays easy to
run and inspect.
"""

from __future__ import annotations

import json
import os
import shutil
import socket
import subprocess
import time
import uuid
from collections.abc import Callable
from datetime import datetime, timezone
from pathlib import Path
from typing import Literal

import fitz
from pydantic import BaseModel, Field

from .errors import PixelRAGError

MAX_PDF_BYTES = 50 * 1024 * 1024
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".png", ".jpg", ".jpeg"}


class DocumentRecord(BaseModel):
    id: str
    name: str
    status: Literal["uploaded", "processing", "ready", "failed"]
    uploaded_at: str
    page_count: int | None = None
    size_bytes: int
    error: str | None = None
    legacy: bool = False
    original_type: str = "pdf"
    normalized_from: str | None = None


class DocumentLibraryState(BaseModel):
    selected_document_id: str | None = None
    documents: list[DocumentRecord] = Field(default_factory=list)


IndexBuilder = Callable[[DocumentRecord, Path, Path], None]


class DocumentLibrary:
    """Filesystem-backed document registry with one PixelRAG index per PDF."""

    LEGACY_ID = "legacy-spm"

    def __init__(
        self,
        root: Path | str,
        *,
        index_builder: IndexBuilder | None = None,
        before_index: Callable[[], None] | None = None,
    ) -> None:
        self.root = Path(root).resolve()
        self.storage_root = self.root / "document-storage"
        self.index_root = self.root / "indexes" / "documents"
        self.registry_path = self.root / "mock-data" / "document-library.json"
        self._index_builder = index_builder or self._build_with_pixelrag
        self._before_index = before_index or (lambda: None)
        self.storage_root.mkdir(parents=True, exist_ok=True)
        self.index_root.mkdir(parents=True, exist_ok=True)

    def _load(self) -> DocumentLibraryState:
        if self.registry_path.exists():
            try:
                state = DocumentLibraryState.model_validate_json(
                    self.registry_path.read_text(encoding="utf-8")
                )
            except (OSError, ValueError):
                state = DocumentLibraryState()
        else:
            state = DocumentLibraryState()
        return self._with_legacy_seed(state)

    def _with_legacy_seed(self, state: DocumentLibraryState) -> DocumentLibraryState:
        legacy_source = self.root / "documents" / "1.pdf"
        legacy_index = self.root / "indexes" / "spm"
        if not legacy_source.exists():
            return state
        if any(item.id == self.LEGACY_ID for item in state.documents):
            return state
        try:
            with fitz.open(legacy_source) as pdf:
                pages = pdf.page_count
        except Exception:
            pages = None
        ready = (legacy_index / "index.faiss").exists() and (legacy_index / "articles.json").exists()
        legacy = DocumentRecord(
            id=self.LEGACY_ID,
            name="Q2 FY2026 Performance Report.pdf",
            status="ready" if ready else "uploaded",
            uploaded_at=datetime.fromtimestamp(legacy_source.stat().st_mtime, timezone.utc).isoformat(),
            page_count=pages,
            size_bytes=legacy_source.stat().st_size,
            legacy=True,
        )
        state.documents.insert(0, legacy)
        if state.selected_document_id is None:
            state.selected_document_id = legacy.id
        return state

    def _save(self, state: DocumentLibraryState) -> None:
        self.registry_path.parent.mkdir(parents=True, exist_ok=True)
        tmp = self.registry_path.with_suffix(".tmp")
        tmp.write_text(state.model_dump_json(indent=2) + "\n", encoding="utf-8")
        tmp.replace(self.registry_path)

    def list_state(self) -> DocumentLibraryState:
        state = self._load()
        # Persist the seed so selection survives subsequent restarts.
        self._save(state)
        return state

    def get(self, document_id: str) -> DocumentRecord:
        state = self._load()
        for item in state.documents:
            if item.id == document_id:
                return item
        raise KeyError(document_id)

    def selected(self) -> DocumentRecord:
        state = self._load()
        if not state.selected_document_id:
            raise KeyError("No document selected")
        return self.get(state.selected_document_id)

    def select(self, document_id: str) -> DocumentRecord:
        state = self._load()
        selected = next((item for item in state.documents if item.id == document_id), None)
        if selected is None:
            raise KeyError(document_id)
        if selected.status != "ready":
            raise ValueError("Only a Ready document can be selected")
        state.selected_document_id = selected.id
        self._save(state)
        return selected

    def source_path(self, document_id: str) -> Path:
        """Normalized PDF consumed by PixelRAG."""
        if document_id == self.LEGACY_ID:
            return self.root / "documents" / "1.pdf"
        return self.storage_root / document_id / "source" / "1.pdf"

    def original_path(self, document_id: str) -> Path:
        record = self.get(document_id)
        if record.legacy:
            return self.source_path(document_id)

        document_dir = self.storage_root / document_id

        # Originals are kept outside the PixelRAG source directory so that
        # PixelRAG sees only the normalized numeric PDF (source/1.pdf).
        matches = list(document_dir.glob("original.*"))
        if matches:
            return matches[0]

        # Backward compatibility for documents created before the storage fix.
        legacy_matches = list((document_dir / "source").glob("original.*"))
        return legacy_matches[0] if legacy_matches else self.source_path(document_id)

    def index_dir(self, document_id: str) -> Path:
        if document_id == self.LEGACY_ID:
            return self.root / "indexes" / "spm"
        return self.index_root / document_id

    def upload(self, filename: str, content: bytes, content_type: str | None = None) -> DocumentRecord:
        safe_name = Path(filename or "uploaded.pdf").name
        extension = Path(safe_name).suffix.casefold()
        if extension not in SUPPORTED_EXTENSIONS:
            raise ValueError("Supported uploads are PDF, DOCX, PPTX, XLSX, PNG, and JPG files")
        if not content:
            raise ValueError("The uploaded document is empty")
        if len(content) > MAX_PDF_BYTES:
            raise ValueError("Document is larger than the 50 MB POC upload limit")

        document_id = f"doc-{uuid.uuid4().hex[:12]}"
        document_dir = self.storage_root / document_id
        source_dir = document_dir / "source"
        source_dir.mkdir(parents=True, exist_ok=False)

        # Preserve the user's original file outside the directory PixelRAG
        # scans. PixelRAG should only see source/1.pdf.
        original = document_dir / f"original{extension}"
        original.write_bytes(content)
        normalized = source_dir / "1.pdf"
        try:
            page_count = self._normalize_to_pdf(original, normalized, extension)
        except Exception:
            shutil.rmtree(self.storage_root / document_id, ignore_errors=True)
            raise

        record = DocumentRecord(
            id=document_id,
            name=safe_name,
            status="uploaded",
            uploaded_at=datetime.now(timezone.utc).isoformat(),
            page_count=page_count,
            size_bytes=len(content),
            original_type=extension.lstrip("."),
            normalized_from=None if extension == ".pdf" else extension.lstrip("."),
        )
        state = self._load()
        state.documents.append(record)
        self._save(state)
        return record

    def _normalize_to_pdf(self, original: Path, normalized: Path, extension: str) -> int:
        if extension == ".pdf":
            data = original.read_bytes()
            if not data.startswith(b"%PDF-"):
                raise ValueError("The uploaded file is not a valid PDF")
            try:
                with fitz.open(original) as pdf:
                    if pdf.page_count < 1:
                        raise ValueError("The uploaded PDF has no pages")
                    pages = pdf.page_count
            except Exception as error:
                raise ValueError("The uploaded file could not be opened as a PDF") from error
            shutil.copy2(original, normalized)
            return pages

        if extension in {".png", ".jpg", ".jpeg"}:
            try:
                pix = fitz.Pixmap(str(original))
                doc = fitz.open()
                page = doc.new_page(width=pix.width, height=pix.height)
                page.insert_image(page.rect, filename=str(original))
                doc.save(normalized)
                doc.close()
                return 1
            except Exception as error:
                raise ValueError("The uploaded image could not be converted to PDF") from error

        lines: list[str] = []
        try:
            if extension == ".docx":
                from docx import Document as WordDocument
                doc = WordDocument(original)
                for para in doc.paragraphs:
                    if para.text.strip():
                        lines.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        lines.append(" | ".join(cell.text.strip() for cell in row.cells))
            elif extension == ".pptx":
                from pptx import Presentation
                prs = Presentation(original)
                for index, slide in enumerate(prs.slides, 1):
                    lines.append(f"--- Slide {index} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text.strip())
            elif extension == ".xlsx":
                from openpyxl import load_workbook
                book = load_workbook(original, read_only=True, data_only=True)
                for sheet in book.worksheets:
                    lines.append(f"--- Worksheet: {sheet.title} ---")
                    for row in sheet.iter_rows(values_only=True):
                        values = ["" if value is None else str(value) for value in row]
                        if any(values):
                            lines.append(" | ".join(values))
        except Exception as error:
            raise ValueError(f"The uploaded {extension.lstrip('.').upper()} file could not be read") from error
        if not lines:
            lines = [f"{original.name}: no extractable text was found."]
        self._text_lines_to_pdf(lines, normalized)
        with fitz.open(normalized) as pdf:
            return pdf.page_count

    @staticmethod
    def _text_lines_to_pdf(lines: list[str], target: Path) -> None:
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)
        y = 50.0
        for raw in lines:
            text = raw.replace("\x00", " ")[:4000]
            chunks = [text[i:i + 95] for i in range(0, len(text), 95)] or [""]
            for chunk in chunks:
                if y > 790:
                    page = doc.new_page(width=595, height=842)
                    y = 50.0
                page.insert_text((45, y), chunk, fontsize=9, fontname="helv")
                y += 13
        doc.save(target)
        doc.close()

    def index(self, document_id: str) -> DocumentRecord:
        record = self.get(document_id)
        if record.legacy:
            if self.index_dir(document_id).joinpath("index.faiss").exists():
                return self._update(document_id, status="ready", error=None)
            raise ValueError("The bundled demo document index is missing. Run `pixelrag index build` first.")

        self._before_index()
        index_dir = self.index_dir(document_id)
        if index_dir.exists():
            shutil.rmtree(index_dir)
        index_dir.parent.mkdir(parents=True, exist_ok=True)
        self._update(document_id, status="processing", error=None)
        try:
            self._index_builder(record, self.source_path(document_id).parent, index_dir)
            required = [index_dir / "index.faiss", index_dir / "articles.json"]
            if not all(path.exists() for path in required):
                raise PixelRAGError("PixelRAG indexing finished without the expected index artifacts")
        except Exception as error:
            message = str(error).strip() or error.__class__.__name__
            self._update(document_id, status="failed", error=message[-2000:])
            raise
        record = self._update(document_id, status="ready", error=None)
        # Newly indexed uploads become the active source automatically.
        self.select(document_id)
        return record

    def _update(self, document_id: str, **changes) -> DocumentRecord:
        state = self._load()
        for index, item in enumerate(state.documents):
            if item.id == document_id:
                updated = item.model_copy(update=changes)
                state.documents[index] = updated
                self._save(state)
                return updated
        raise KeyError(document_id)

    def _build_with_pixelrag(self, record: DocumentRecord, source_dir: Path, output_dir: Path) -> None:
        executable = shutil.which("pixelrag")
        if not executable:
            raise PixelRAGError("The `pixelrag` command was not found. Activate the project .venv first.")
        work_dir = self.storage_root / record.id / "index-work"
        work_dir.mkdir(parents=True, exist_ok=True)
        config = work_dir / "pixelrag.yaml"
        config.write_text(
            "\n".join([
                "source:",
                "  type: local",
                f"  path: {json.dumps(str(source_dir))}",
                "",
                "embed:",
                "  model: Qwen/Qwen3-VL-Embedding-2B",
                "  device: auto",
                "",
                f"output: {json.dumps(str(output_dir))}",
                "",
            ]),
            encoding="utf-8",
        )
        log_path = self.storage_root / record.id / "index.log"
        with log_path.open("w", encoding="utf-8") as log:
            process = subprocess.run(
                [executable, "index", "build"],
                cwd=work_dir,
                stdout=log,
                stderr=subprocess.STDOUT,
                text=True,
                timeout=int(os.getenv("PIXELRAG_INDEX_TIMEOUT", "1800")),
                check=False,
            )
        if process.returncode != 0:
            try:
                lines = log_path.read_text(encoding="utf-8", errors="replace").splitlines()
                tail = "\n".join(lines[-20:])
            except OSError:
                tail = ""
            raise PixelRAGError(
                f"PixelRAG indexing failed for {record.name}. "
                f"See {log_path}." + (f"\n{tail}" if tail else "")
            )


class PixelRAGProcessManager:
    """Runs one local PixelRAG search server for the currently selected index.

    A single process is intentional: the Qwen visual embedding model is large and the
    target laptop has 8 GB VRAM. Starting one server per document would waste VRAM.
    """

    def __init__(self, root: Path | str, port: int | None = None) -> None:
        self.root = Path(root).resolve()
        self.port = port or int(os.getenv("PIXELRAG_MANAGED_PORT", "30002"))
        self.process: subprocess.Popen | None = None
        self.document_id: str | None = None
        self.log_handle = None

    @property
    def base_url(self) -> str:
        return f"http://127.0.0.1:{self.port}"

    def stop(self) -> None:
        if self.process and self.process.poll() is None:
            self.process.terminate()
            try:
                self.process.wait(timeout=10)
            except subprocess.TimeoutExpired:
                self.process.kill()
                self.process.wait(timeout=5)
        self.process = None
        self.document_id = None
        if self.log_handle:
            self.log_handle.close()
            self.log_handle = None

    def ensure(self, document: DocumentRecord, index_dir: Path) -> str:
        if self.process and self.process.poll() is None and self.document_id == document.id:
            return self.base_url
        self.stop()
        executable = shutil.which("pixelrag")
        if not executable:
            raise PixelRAGError("The `pixelrag` command was not found. Activate the project .venv first.")
        articles = index_dir / "articles.json"
        if not (index_dir / "index.faiss").exists() or not articles.exists():
            raise PixelRAGError(f"PixelRAG index is not ready for {document.name}")
        log_path = self.root / "document-storage" / "pixelrag-active.log"
        log_path.parent.mkdir(parents=True, exist_ok=True)
        self.log_handle = log_path.open("w", encoding="utf-8")
        self.process = subprocess.Popen(
            [
                executable,
                "serve",
                "--index-dir", str(index_dir),
                "--articles-json", str(articles),
                "--port", str(self.port),
            ],
            cwd=self.root,
            stdout=self.log_handle,
            stderr=subprocess.STDOUT,
            text=True,
        )
        self.document_id = document.id
        deadline = time.monotonic() + int(os.getenv("PIXELRAG_SERVE_TIMEOUT", "180"))
        while time.monotonic() < deadline:
            if self.process.poll() is not None:
                break
            try:
                with socket.create_connection(("127.0.0.1", self.port), timeout=0.5):
                    return self.base_url
            except OSError:
                time.sleep(0.5)
        code = self.process.poll()
        self.stop()
        try:
            tail = "\n".join(log_path.read_text(encoding="utf-8", errors="replace").splitlines()[-20:])
        except OSError:
            tail = ""
        raise PixelRAGError(
            f"PixelRAG search service did not start for {document.name}"
            + (f" (exit code {code})" if code is not None else "")
            + (f".\n{tail}" if tail else ".")
        )
