"""Office-aware document normalization for PixelRAG.

LibreOffice is preferred when available because it preserves the original visual
layout.  The pure-Python fallbacks keep meaningful document boundaries so the
resulting PixelRAG evidence remains traceable even on lightweight hosts where
LibreOffice is not installed.
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

import fitz

from .documents import DocumentLibrary


class VisualDocumentLibrary(DocumentLibrary):
    """Prefer true Office-to-PDF rendering, with structure-preserving fallbacks."""

    OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}

    def _normalize_to_pdf(self, original: Path, normalized: Path, extension: str) -> int:
        if extension not in self.OFFICE_EXTENSIONS:
            return super()._normalize_to_pdf(original, normalized, extension)

        if self._try_libreoffice(original, normalized):
            with fitz.open(normalized) as pdf:
                if pdf.page_count < 1:
                    raise ValueError("Converted Office document has no pages")
                return pdf.page_count

        if extension == ".docx":
            return self._normalize_docx_fallback(original, normalized)
        if extension == ".pptx":
            return self._normalize_pptx_fallback(original, normalized)
        if extension == ".xlsx":
            return self._normalize_xlsx_fallback(original, normalized)

        return super()._normalize_to_pdf(original, normalized, extension)

    @staticmethod
    def _try_libreoffice(original: Path, normalized: Path) -> bool:
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return False

        try:
            with tempfile.TemporaryDirectory(prefix="pixelrag-office-") as temporary:
                output = Path(temporary)
                process = subprocess.run(
                    [
                        soffice,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        str(output),
                        str(original),
                    ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.STDOUT,
                    text=True,
                    timeout=180,
                    check=False,
                )
                generated = output / f"{original.stem}.pdf"
                if process.returncode != 0 or not generated.is_file():
                    return False

                shutil.copy2(generated, normalized)
                return True
        except (OSError, subprocess.SubprocessError):
            return False

    @staticmethod
    def _paragraph_tokens(paragraph) -> list[str | None]:
        """Return DOCX paragraph text with None tokens representing page breaks."""
        tokens: list[str | None] = []
        buffer: list[str] = []

        def flush() -> None:
            value = "".join(buffer).strip()
            if value:
                tokens.append(value)
            buffer.clear()

        for element in paragraph._p.iter():
            local_name = element.tag.rsplit("}", 1)[-1]

            if local_name == "t":
                buffer.append(element.text or "")
            elif local_name == "tab":
                buffer.append("\t")
            elif local_name == "br":
                break_type = next(
                    (value for key, value in element.attrib.items() if key.endswith("}type")),
                    None,
                )
                if break_type == "page":
                    flush()
                    tokens.append(None)
                else:
                    buffer.append(" ")
            elif local_name == "lastRenderedPageBreak":
                flush()
                tokens.append(None)

        flush()
        return tokens

    @staticmethod
    def _normalize_docx_fallback(original: Path, normalized: Path) -> int:
        """Preserve Word document order and explicit/rendered page boundaries."""
        try:
            from docx import Document as WordDocument
            from docx.table import Table
            from docx.text.paragraph import Paragraph

            document = WordDocument(original)
            pages: list[list[str]] = [[]]

            for child in document.element.body.iterchildren():
                if child.tag.endswith("}p"):
                    paragraph = Paragraph(child, document)

                    if paragraph.paragraph_format.page_break_before is True and pages[-1]:
                        pages.append([])

                    for token in VisualDocumentLibrary._paragraph_tokens(paragraph):
                        if token is None:
                            pages.append([])
                        elif token:
                            pages[-1].append(token)

                elif child.tag.endswith("}tbl"):
                    table = Table(child, document)
                    for row in table.rows:
                        text = " | ".join(cell.text.strip() for cell in row.cells)
                        if text.strip(" |"):
                            pages[-1].append(text)

            if not any(page for page in pages):
                pages = [[f"{original.name}: no extractable text was found."]]

            VisualDocumentLibrary._text_pages_to_pdf(pages, normalized)
            return VisualDocumentLibrary._validated_page_count(normalized, "DOCX")
        except Exception as error:
            raise ValueError("The uploaded DOCX file could not be read") from error

    @staticmethod
    def _normalize_pptx_fallback(original: Path, normalized: Path) -> int:
        """Preserve one logical PDF page per PowerPoint slide."""
        try:
            from pptx import Presentation

            presentation = Presentation(original)
            pages: list[list[str]] = []

            for slide_number, slide in enumerate(presentation.slides, start=1):
                lines = [f"Slide {slide_number}"]

                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text = (getattr(shape, "text", "") or "").strip()
                        if text:
                            lines.extend(part.strip() for part in text.splitlines() if part.strip())

                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            text = " | ".join(cell.text.strip() for cell in row.cells)
                            if text.strip(" |"):
                                lines.append(text)

                pages.append(lines)

            if not pages:
                pages = [[f"{original.name}: no slides were found."]]

            VisualDocumentLibrary._text_pages_to_pdf(pages, normalized)
            return VisualDocumentLibrary._validated_page_count(normalized, "PPTX")
        except Exception as error:
            raise ValueError("The uploaded PPTX file could not be read") from error

    @staticmethod
    def _normalize_xlsx_fallback(original: Path, normalized: Path) -> int:
        """Preserve worksheet boundaries, allowing large sheets to overflow naturally."""
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(original, read_only=True, data_only=True)
            pages: list[list[str]] = []

            try:
                for sheet in workbook.worksheets:
                    lines = [f"Worksheet: {sheet.title}"]
                    for row in sheet.iter_rows(values_only=True):
                        values = ["" if value is None else str(value) for value in row]
                        if any(value.strip() for value in values):
                            lines.append(" | ".join(values))
                    pages.append(lines)
            finally:
                workbook.close()

            if not pages:
                pages = [[f"{original.name}: no worksheets were found."]]

            VisualDocumentLibrary._text_pages_to_pdf(pages, normalized)
            return VisualDocumentLibrary._validated_page_count(normalized, "XLSX")
        except Exception as error:
            raise ValueError("The uploaded XLSX file could not be read") from error

    @staticmethod
    def _validated_page_count(path: Path, label: str) -> int:
        with fitz.open(path) as pdf:
            if pdf.page_count < 1:
                raise ValueError(f"Converted {label} document has no pages")
            return pdf.page_count

    @staticmethod
    def _text_pages_to_pdf(pages: list[list[str]], target: Path) -> None:
        """Create one PDF page per logical source page/sheet/slide plus overflow pages."""
        pdf = fitz.open()

        for source_page in pages:
            page = pdf.new_page(width=595, height=842)
            y = 50.0

            for raw in source_page:
                text = raw.replace("\x00", " ")[:4000]
                chunks = [text[index:index + 95] for index in range(0, len(text), 95)] or [""]

                for chunk in chunks:
                    if y > 790:
                        page = pdf.new_page(width=595, height=842)
                        y = 50.0
                    page.insert_text((45, y), chunk, fontsize=9, fontname="helv")
                    y += 13

        pdf.save(target)
        pdf.close()
